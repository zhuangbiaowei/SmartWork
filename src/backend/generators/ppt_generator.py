from typing import Dict, Any
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None
    Inches = None
    Pt = None
    PP_ALIGN = None

from .document_generator import DocumentGenerator


class PowerPointGenerator(DocumentGenerator):
    """
    PowerPoint document generator using python-pptx.

    Supports slides with text, shapes, and formatting.
    """

    def __init__(self):
        self.prs = None

    async def generate(
        self, data: Dict[str, Any], output_path: str, **kwargs: Any
    ) -> bool:
        """
        Generate PowerPoint presentation from data.

        Args:
            data: Dictionary with:
                - title: Presentation title
                - slides: List of slides (each slide is a dict)
            output_path: Path to save PowerPoint file

        Returns:
            True if successful, False otherwise
        """
        if Presentation is None:
            print("python-pptx not installed. Install with: pip install python-pptx")
            return False

        try:
            self.prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            title_slide = prs.slides.add_slide(0, 0, prs.slide_width, prs.slide_height)
            title_box = title_slide.textbox
            title_box.text = data.get("title", "Presentation")
            title_box.text_frame.paragraph_format.alignment = PP_ALIGN.CENTER
            title_box.font.size = Pt(44)
            title_box.font.bold = True

            if "slides" in data:
                for idx, slide_data in enumerate(data["slides"], 1):
                    slide = prs.slides.add_slide(
                        idx, 0, prs.slide_width, prs.slide_height
                    )

                    if "title" in slide_data:
                        title = slide.shapes.title
                        title.text = slide_data["title"]
                        title.text_frame.paragraph_format.alignment = PP_ALIGN.CENTER
                        title.font.size = Pt(32)
                        title.font.bold = True

                    if "content" in slide_data:
                        content_box = slide.shapes.textbox
                        content_box.text = slide_data["content"]
                        content_box.text_frame.word_wrap = True
                        content_box.text_frame.paragraph_format.alignment = (
                            PP_ALIGN.LEFT
                        )
                        content_box.font.size = Pt(18)

                    if "points" in slide_data:
                        points_text = slide.shapes.textbox
                        points_text.text = "• " + "\n• ".join(
                            f"  {point}" for point in slide_data["points"]
                        )
                        points_text.left = Inches(0.5)
                        points_text.top = Inches(1.5)
                        points_text.font.size = Pt(16)

            file_path = Path(output_path)
            file_path.parent.mkdir(parents=True, exist_ok=True)

            self.prs.save(file_path)
            return True

        except Exception as e:
            print(f"Failed to generate PowerPoint: {e}")
            return False
